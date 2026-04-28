"""
Generates presentation.pptx — a full slide deck summarizing the multi-agent debate experiment.
Run with: python3 build_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import json, glob, numpy as np

# ── Palette ────────────────────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x0F, 0x17, 0x2A)   # deep navy
BG_CARD   = RGBColor(0x1A, 0x26, 0x40)   # slightly lighter navy
ACCENT    = RGBColor(0x4C, 0x9B, 0xE8)   # blue
ACCENT2   = RGBColor(0x6D, 0xBF, 0x6D)   # green
WARN      = RGBColor(0xE8, 0x5C, 0x5C)   # red
GOLD      = RGBColor(0xF5, 0xA6, 0x23)   # orange/gold
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xB0, 0xC4, 0xDE)   # muted blue-white

AGENT_COLORS = {
    "Alice": ACCENT,
    "Bob":   WARN,
    "Carol": ACCENT2,
    "David": GOLD,
}

W = Inches(13.33)
H = Inches(7.5)

# ── Helpers ────────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txBox

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_slide(prs, layout_idx=6):
    layout = prs.slide_layouts[layout_idx]   # blank
    return prs.slides.add_slide(layout)

def card(slide, x, y, w, h, title, body, title_color=ACCENT, body_size=14):
    add_rect(slide, x, y, w, h, BG_CARD)
    add_text(slide, title, x + Inches(0.15), y + Inches(0.12),
             w - Inches(0.3), Inches(0.4), size=13, bold=True, color=title_color)
    add_text(slide, body,  x + Inches(0.15), y + Inches(0.5),
             w - Inches(0.3), h - Inches(0.6), size=body_size, color=LIGHT)

# ── Load trial data ─────────────────────────────────────────────────────────────
def load_trials():
    paths = sorted(glob.glob("trial_*.json"))
    return [json.load(open(p, encoding="utf-8")) for p in paths]

trials = load_trials()

# Final beliefs per agent (last trial as representative)
def final_belief(agent_name, trial_idx=-1):
    for ag in trials[trial_idx]["agents"]:
        if ag["name"] == agent_name:
            return ag["final_belief"]
    return ""

# ── Build deck ─────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ── SLIDE 1: Title ─────────────────────────────────────────────────────────────
sl = add_slide(prs)
set_bg(sl, BG_DARK)
add_rect(sl, 0, Inches(2.6), W, Inches(0.5), ACCENT)   # accent bar

add_text(sl, "Multi-Agent Debate Simulation",
         Inches(0.8), Inches(1.1), Inches(11), Inches(1.2),
         size=44, bold=True, align=PP_ALIGN.CENTER)

add_text(sl, "Can AI characters change each other's minds?",
         Inches(0.8), Inches(2.2), Inches(11), Inches(0.5),
         size=22, color=LIGHT, align=PP_ALIGN.CENTER)

add_text(sl, "Topic: Should students use AI tools in education?\n5 Trials  ·  5 Rounds each  ·  4 Characters",
         Inches(0.8), Inches(3.3), Inches(11), Inches(0.8),
         size=17, color=LIGHT, align=PP_ALIGN.CENTER)

# character chips
chip_labels = ["Alice — Optimistic", "Bob — Skeptical", "Carol — Mediator", "David — Cautious"]
chip_colors = [ACCENT, WARN, ACCENT2, GOLD]
chip_w = Inches(2.8)
start_x = Inches(0.65)
for i, (lbl, col) in enumerate(zip(chip_labels, chip_colors)):
    add_rect(sl, start_x + i * (chip_w + Inches(0.2)), Inches(4.4), chip_w, Inches(0.55), col)
    add_text(sl, lbl, start_x + i * (chip_w + Inches(0.2)) + Inches(0.1), Inches(4.45),
             chip_w - Inches(0.2), Inches(0.45), size=13, bold=True, align=PP_ALIGN.CENTER)

add_text(sl, "Powered by Claude claude-sonnet-4-6  ·  Anthropic",
         Inches(0), Inches(7.0), W, Inches(0.35),
         size=11, color=RGBColor(0x55, 0x6B, 0x8A), align=PP_ALIGN.CENTER)

# ── SLIDE 2: How It Works ──────────────────────────────────────────────────────
sl = add_slide(prs)
set_bg(sl, BG_DARK)
add_text(sl, "How the Simulation Works", Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
         size=30, bold=True)
add_rect(sl, Inches(0.6), Inches(0.95), Inches(4.5), Pt(3), ACCENT)

steps = [
    ("1  Set up characters",
     "Four AI personas are created with different starting beliefs and goals about AI in education."),
    ("2  Run a round",
     "Each character reads what others said, thinks about it, updates their own belief, and writes a new message."),
    ("3  Memory",
     "Important moments are saved to each character's memory so they can reference them in later rounds."),
    ("4  Repeat × 5 rounds",
     "The debate continues for 5 rounds per trial, giving characters time to shift or hold their ground."),
    ("5  Analyze",
     "After all 5 trials finish, charts and summaries are generated to show what changed and why."),
]
col_w = Inches(2.35)
for i, (title, body) in enumerate(steps):
    cx = Inches(0.35) + i * (col_w + Inches(0.2))
    card(sl, cx, Inches(1.2), col_w, Inches(5.7), title, body, title_color=ACCENT, body_size=13)

add_rect(sl, 0, Inches(7.15), W, Inches(0.35), BG_CARD)
add_text(sl, "Each trial is fully independent — characters start fresh with no memory of previous trials.",
         Inches(0.5), Inches(7.15), Inches(12), Inches(0.35),
         size=12, color=LIGHT, align=PP_ALIGN.CENTER)

# ── SLIDE 3: The Characters ────────────────────────────────────────────────────
sl = add_slide(prs)
set_bg(sl, BG_DARK)
add_text(sl, "Meet the Characters", Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
         size=30, bold=True)
add_rect(sl, Inches(0.6), Inches(0.95), Inches(4.5), Pt(3), ACCENT)

chars = [
    ("Alice", ACCENT, "Optimistic",
     "Starting belief: AI tools can support learning and creativity, but I'm still figuring out where the right limits are.",
     "Goal: Promote the benefits of AI in education while staying open to criticism."),
    ("Bob", WARN, "Skeptical",
     "Starting belief: AI tools carry real risks of dependency and misuse, though I can imagine responsible use cases.",
     "Goal: Surface the risks and blind spots in overly optimistic views."),
    ("Carol", ACCENT2, "Mediator",
     "Starting belief: AI tools are useful, but schools need clear boundaries.",
     "Goal: Mediate between both sides and push the group toward balanced guidelines."),
    ("David", GOLD, "Cautious Supporter",
     "Starting belief: AI tools can help students, but only when used responsibly.",
     "Goal: Support useful adoption of AI while emphasizing moderation."),
]

card_w = Inches(2.9)
for i, (name, color, role, belief, goal) in enumerate(chars):
    cx = Inches(0.3) + i * (card_w + Inches(0.22))
    add_rect(sl, cx, Inches(1.2), card_w, Inches(5.8), BG_CARD)
    add_rect(sl, cx, Inches(1.2), card_w, Inches(0.55), color)
    add_text(sl, name, cx + Inches(0.15), Inches(1.22), card_w - Inches(0.3), Inches(0.35),
             size=18, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl, role, cx + Inches(0.15), Inches(1.65), card_w - Inches(0.3), Inches(0.32),
             size=12, color=color, align=PP_ALIGN.CENTER)
    add_text(sl, belief, cx + Inches(0.15), Inches(2.05), card_w - Inches(0.3), Inches(2.2),
             size=12, color=LIGHT)
    add_text(sl, goal, cx + Inches(0.15), Inches(4.4), card_w - Inches(0.3), Inches(2.4),
             size=12, color=WHITE)

# ── SLIDE 4: Stance Evolution chart ───────────────────────────────────────────
sl = add_slide(prs)
set_bg(sl, BG_DARK)
add_text(sl, "Chart 1 — Did Their Opinions Change Each Round?", Inches(0.6), Inches(0.3),
         Inches(12), Inches(0.6), size=28, bold=True)
add_rect(sl, Inches(0.6), Inches(0.92), Inches(5), Pt(3), ACCENT)
add_text(sl,
    "Each line is one of the 5 debate trials. The Y-axis shows whether the character was supportive, "
    "balanced, or skeptical in that round. Overlapping lines mean the character behaved consistently "
    "no matter the trial.",
    Inches(0.6), Inches(1.05), Inches(7.5), Inches(1.0), size=13, color=LIGHT)
sl.shapes.add_picture("stance_evolution.png", Inches(0.4), Inches(2.0), Inches(12.5), Inches(5.2))

# ── SLIDE 5: Convergence Curve ────────────────────────────────────────────────
sl = add_slide(prs)
set_bg(sl, BG_DARK)
add_text(sl, "Chart 2 — Did the Group Start Agreeing Over Time?", Inches(0.6), Inches(0.3),
         Inches(12), Inches(0.6), size=28, bold=True)
add_rect(sl, Inches(0.6), Inches(0.92), Inches(5), Pt(3), ACCENT2)
add_text(sl,
    "Lower values = the 4 characters' opinions are closer together. A downward curve means the debate "
    "brought them toward agreement. A flat or rising curve means they stayed divided.",
    Inches(0.6), Inches(1.05), Inches(7.5), Inches(0.9), size=13, color=LIGHT)
sl.shapes.add_picture("convergence_curve.png", Inches(1.5), Inches(2.0), Inches(10.2), Inches(5.1))

# ── SLIDE 6: Influence Network ────────────────────────────────────────────────
sl = add_slide(prs)
set_bg(sl, BG_DARK)
add_text(sl, "Chart 3 — Who Was Convincing Whom?", Inches(0.6), Inches(0.3),
         Inches(12), Inches(0.6), size=28, bold=True)
add_rect(sl, Inches(0.6), Inches(0.92), Inches(4), Pt(3), GOLD)
add_text(sl,
    "Arrows point from the influencer to the person they moved. Thicker arrows = that influence "
    "happened more often across all 5 trials. The number on each arrow is the total count.",
    Inches(0.6), Inches(1.05), Inches(6.5), Inches(0.9), size=13, color=LIGHT)
sl.shapes.add_picture("influence_network.png", Inches(2.5), Inches(1.4), Inches(8.0), Inches(5.9))

# ── SLIDE 7: Belief Drift ─────────────────────────────────────────────────────
sl = add_slide(prs)
set_bg(sl, BG_DARK)
add_text(sl, "Chart 4 — How Much Did Each Character Actually Change?", Inches(0.6), Inches(0.3),
         Inches(12), Inches(0.6), size=26, bold=True)
add_rect(sl, Inches(0.6), Inches(0.92), Inches(5), Pt(3), WARN)
add_text(sl,
    "0 = no change from starting belief. 1 = completely shifted. The shaded band shows how consistent "
    "this was — a wide band means some trials moved them a lot, others barely at all.",
    Inches(0.6), Inches(1.05), Inches(7.5), Inches(0.9), size=13, color=LIGHT)
sl.shapes.add_picture("belief_drift.png", Inches(0.4), Inches(2.0), Inches(12.5), Inches(5.2))

# ── SLIDE 8: Venn Diagram ─────────────────────────────────────────────────────
sl = add_slide(prs)
set_bg(sl, BG_DARK)
add_text(sl, "Chart 5 — Which Topics Kept Coming Up?", Inches(0.6), Inches(0.3),
         Inches(12), Inches(0.6), size=28, bold=True)
add_rect(sl, Inches(0.6), Inches(0.92), Inches(4), Pt(3), ACCENT)
add_text(sl,
    "Each circle pair compares two consecutive trials. The overlapping center = topics both sessions "
    "discussed. The outer edges = topics unique to just one trial.",
    Inches(0.6), Inches(1.05), Inches(7), Inches(0.9), size=13, color=LIGHT)
sl.shapes.add_picture("venn_diagram.png", Inches(0.3), Inches(1.9), Inches(12.7), Inches(5.4))

# ── SLIDE 9: Key Findings ─────────────────────────────────────────────────────
sl = add_slide(prs)
set_bg(sl, BG_DARK)
add_text(sl, "Key Findings", Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
         size=30, bold=True)
add_rect(sl, Inches(0.6), Inches(0.92), Inches(2.5), Pt(3), ACCENT)

findings = [
    (ACCENT,  "Broad agreement emerged",
     "Across all 5 trials, every character ended up agreeing that AI in schools requires careful guidelines "
     "and a strong focus on preserving students' critical thinking."),
    (WARN,    "Bob was the hardest to move",
     "Bob (the skeptic) held his ground most consistently — his belief that AI poses real cognitive risks "
     "stayed firm in nearly every trial."),
    (ACCENT2, "Carol and David drove convergence",
     "Carol (mediator) and David (cautious supporter) were the characters most likely to shift toward "
     "common ground, helping pull the group together."),
    (GOLD,    "Critical thinking was universal",
     "The theme of 'critical thinking' appeared in all 5 trials without exception — it was the "
     "single most consistent topic the simulation gravitated toward."),
    (ACCENT,  "Over-reliance was the core fear",
     "Even the supportive characters (Alice, Carol) ended up worried about students becoming too dependent "
     "on AI — the risk of over-reliance outlasted any initial optimism."),
]

card_w = Inches(2.35)
for i, (color, title, body) in enumerate(findings):
    cx = Inches(0.3) + i * (card_w + Inches(0.22))
    card(sl, cx, Inches(1.2), card_w, Inches(5.9), title, body, title_color=color, body_size=13)

# ── SLIDE 10: Final Beliefs ────────────────────────────────────────────────────
sl = add_slide(prs)
set_bg(sl, BG_DARK)
add_text(sl, "Where Each Character Ended Up", Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
         size=30, bold=True)
add_rect(sl, Inches(0.6), Inches(0.92), Inches(4), Pt(3), ACCENT)
add_text(sl, "Final beliefs from the last trial:", Inches(0.6), Inches(1.05), Inches(12), Inches(0.35),
         size=13, color=LIGHT)

for i, (name, color) in enumerate(AGENT_COLORS.items()):
    belief = final_belief(name)
    cx = Inches(0.3) + (i % 2) * Inches(6.4)
    cy = Inches(1.55) + (i // 2) * Inches(2.8)
    add_rect(sl, cx, cy, Inches(6.1), Inches(2.6), BG_CARD)
    add_rect(sl, cx, cy, Inches(0.18), Inches(2.6), color)
    add_text(sl, name, cx + Inches(0.28), cy + Inches(0.12),
             Inches(5.7), Inches(0.38), size=15, bold=True, color=color)
    add_text(sl, f'"{belief}"', cx + Inches(0.28), cy + Inches(0.52),
             Inches(5.7), Inches(1.95), size=12, color=LIGHT)

# ── SLIDE 11: Consensus ────────────────────────────────────────────────────────
sl = add_slide(prs)
set_bg(sl, BG_DARK)
add_text(sl, "Overall Consensus", Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
         size=30, bold=True)
add_rect(sl, Inches(0.6), Inches(0.92), Inches(3), Pt(3), ACCENT2)

consensus = (
    "Across all 5 trials, a clear consensus emerged: AI tools can be valuable in education, "
    "but only when implemented with strong safeguards for critical thinking.\n\n"
    "All four characters — even the initially optimistic ones — converged on the view that "
    "students must be actively taught to evaluate and question AI-generated content, rather "
    "than passively accepting it.\n\n"
    "The biggest remaining disagreement was about degree of risk: Bob and David stayed more "
    "alarmed about long-term cognitive harm, while Alice and Carol leaned toward cautious optimism."
)
add_rect(sl, Inches(0.6), Inches(1.2), Inches(12.1), Inches(4.8), BG_CARD)
add_text(sl, consensus, Inches(0.9), Inches(1.4), Inches(11.5), Inches(4.4),
         size=17, color=WHITE)

bullets = ["✓  Critical thinking must be actively protected",
           "✓  Clear guidelines are non-negotiable",
           "✓  AI as a tool — not a replacement for thinking",
           "✓  Over-reliance is the central risk to manage"]
for i, b in enumerate(bullets):
    add_rect(sl, Inches(1.1), Inches(6.15) + i * Inches(0.32), Inches(11), Inches(0.28),
             BG_DARK)
    add_text(sl, b, Inches(1.2), Inches(6.15) + i * Inches(0.32), Inches(11), Inches(0.28),
             size=13, color=ACCENT2, bold=True)

# ── SLIDE 12: What This Proves ────────────────────────────────────────────────
sl = add_slide(prs)
set_bg(sl, BG_DARK)
add_text(sl, "What This Experiment Demonstrates", Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
         size=28, bold=True)
add_rect(sl, Inches(0.6), Inches(0.92), Inches(5), Pt(3), ACCENT)

points = [
    (ACCENT, "AI agents can hold genuine opinions",
     "Each character maintained a consistent worldview across multiple trials, not just outputting generic text."),
    (ACCENT2, "Beliefs can change through conversation",
     "Characters demonstrably updated their stated beliefs in response to arguments made by others — not randomly, but logically."),
    (GOLD, "Memory makes debates richer",
     "With persistent memory, characters referenced earlier points later in the debate, creating realistic back-and-forth."),
    (WARN, "Consensus is possible — but hard",
     "Even with only 4 characters, full agreement was never reached. Some tension always remained, just like real debates."),
    (LIGHT, "This is a research tool",
     "This simulation can stress-test any policy question, stakeholder position, or design decision by letting AI personas argue it out."),
]

card_w = Inches(2.35)
for i, (color, title, body) in enumerate(points):
    cx = Inches(0.3) + i * (card_w + Inches(0.22))
    card(sl, cx, Inches(1.2), card_w, Inches(5.9), title, body, title_color=color, body_size=13)

# ── SLIDE 13: End ─────────────────────────────────────────────────────────────
sl = add_slide(prs)
set_bg(sl, BG_DARK)
add_rect(sl, 0, Inches(2.8), W, Inches(0.5), ACCENT)
add_text(sl, "Thank You", Inches(0.5), Inches(1.1), Inches(12.3), Inches(1.4),
         size=54, bold=True, align=PP_ALIGN.CENTER)
add_text(sl, "5 trials  ·  5 rounds each  ·  4 characters  ·  1 conclusion",
         Inches(0.5), Inches(2.2), Inches(12.3), Inches(0.5),
         size=18, color=LIGHT, align=PP_ALIGN.CENTER)
add_text(sl, "AI tools can help — but only with guardrails.",
         Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.8),
         size=28, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)
add_text(sl, "Built with Claude claude-sonnet-4-6 + Python  ·  github.com/christopherwilliams",
         Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4),
         size=11, color=RGBColor(0x55, 0x6B, 0x8A), align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
prs.save("presentation.pptx")
print("Saved → presentation.pptx")
print(f"  {len(prs.slides)} slides")
